from django.shortcuts import render,redirect,get_object_or_404
from django.contrib.auth.models import User
from django.contrib.auth.forms import authenticate
from django.contrib import messages
from django.contrib.auth import login,logout
import os
from docx import Document
import PyPDF2
from PyPDF2 import PdfReader
from pptx import Presentation
import base64
from docx import Document as DocxDocument
from django.http import HttpResponse
from django.contrib.auth.forms import PasswordChangeForm
from django.contrib.auth import update_session_auth_hash
from .forms import SignUpForm,LogInForm
from .models import AboutImageStore,AboutVideoStore
from .forms import AboutVideoForm,AboutImageForm
from .forms import ServiceForm
from .models import ServiceStore
from .models import NewsImageStore,NewsVideoStore
from .forms import NewsImageForm,NewsVideoForm
from .models import SlideStore
from .forms import SlideForm
from .forms import DocumentForm
from .models import DocumentStore
from .models import TouristImageStore,TouristVideoStore
from .forms import TouristImageForm,TouristVideoForm
from .forms import AboutImageRename,AboutVideoRename
from .forms import NewsVideoRename,NewsImageRename
from .forms import ServiceRename,SlideRename,DocumentRename
from .forms import TouristImageRename,TouristVideoRename

def owner(request):
    return render(request, 'owner.html')


def sign_page(request):
    count_number = User.objects.count()
    if count_number>=2:
        return redirect('public')

    if request.method == 'POST':
        form = SignUpForm(request.POST)
        if form.is_valid():
            user = form.save()
            return redirect('mis')
    else:
        form = SignUpForm()
    return render(request,'signup.html',{'form':form})

def login_page(request):
    if request.method == 'POST':
        form = LogInForm(request.POST)
        username = request.POST.get('username')
        password = request.POST.get('password')
        user = authenticate(username=username, password=password)
        if user is None:
            messages.error(request,'invalid username or password')
        else:
            login(request, user)
            return redirect('mis')
    else:
        form = LogInForm()
    return render(request,'login.html',{'form':form})

def logout_page(request):
    logout(request)
    return redirect('MIS-login')

def change_password(request):
    if request.method == 'POST':
        form = PasswordChangeForm(request.user, request.POST)
        if form.is_valid():
            user = form.save()
            update_session_auth_hash(request,user)
            return redirect('mis')
    else:
        form = PasswordChangeForm(request.user)
    return render(request, 'change_password.html',{'form':form})


def about_image_page(request):
    if request.method == 'POST':
        form = AboutImageForm(request.POST,request.FILES)
        if form.is_valid():
            doc =form.save(commit=False)
            doc.user = request.user
            doc.save()
            messages.success(request,'your image have been uploaded')
            return redirect('mis')
    else:
        form = AboutImageForm()
    return render(request,'about_image.html',{'form':form})

def about_image_delete(request, pk):
    image = get_object_or_404(AboutImageStore, pk=pk, user=request.user)
    image.about_image.delete()
    image.delete()
    return redirect('about_owner')

def about_video_page(request):
    if request.method == 'POST':
        form = AboutVideoForm(request.POST,request.FILES)
        if form.is_valid():
            doc = form.save(commit=False)
            doc.user = request.user
            doc.save()
            messages.success(request,'your video have been uploaded')
            return redirect('mis')

    else:
        form = AboutVideoForm()
    return render(request,'about_video.html',{'form':form})

def about_video__delete(request, pk):
    video = get_object_or_404(AboutVideoStore,pk=pk, user=request.user)
    video.about_video.delete()
    video.delete()
    return redirect('about_owner')


def about_list(request):
    context = {'about_image':AboutImageStore.objects.all(),
               'about_video':AboutVideoStore.objects.all()}
    return render(request,'about.html',{'context':context})

def about_owner(request):
    context = {'about_image':AboutImageStore.objects.filter(user=request.user),
               'about_video':AboutVideoStore.objects.filter(user=request.user)}
    return render(request,'owner/about.html',{'context':context})

def about_video_rename(request, pk):
    video = get_object_or_404(AboutVideoStore, pk=pk)

    if request.method =='POST':
        form = AboutVideoRename(request.POST, instance=video)
        if form.is_valid():
            doc = form.save(commit=False)
            doc.save()
            messages.success(request, 'video title reset successfully')
            return redirect('about_owner')
    else:
        form = AboutVideoRename(instance=video)
    return render(request,'rename/about_video.html',{'form':form})

def about_image_rename(request, pk):
    image = get_object_or_404(AboutImageStore, pk=pk)

    if request.method == 'POST':
        form = AboutImageRename(request.POST, instance=image)
        if form.is_valid():
            doc = form.save(commit=False)
            doc.save()
            messages.success(request, 'image reset successfully')
            return redirect('about_owner')
    else:
        form = AboutImageRename()
    return render(request,'rename/about_image.html',{"form":form})



#def public(request):
#    return render(request,'home.html')

def mis(request):
    return render(request,'admin.html')


#SERVICE INFORMATION
def service_page(request):
    if request.method == 'POST':
        form = ServiceForm(request.POST)
        if form.is_valid():
            doc = form.save(commit=False)
            doc.user = request.user
            doc.save()
            messages.success(request,'your text have been upload')
            return redirect('mis')
    else:
        form = ServiceForm()
    return render(request,'service_upload.html',{'form':form})

def service_list(request):
    file = ServiceStore.objects.all()
    return render(request,'service.html',{'file':file})

def service_delete(request, pk):
    file = get_object_or_404(ServiceStore,id=id, user=request.user)
    file.text.delete()
    file.delete()
    return redirect('service_owner')

def service_owner(request):
    file = ServiceStore.objects.filter(user=request.user)
    return render(request,'owner/service.html',{'file':file})

def service_text_rename(request,pk):
    text = get_object_or_404(ServiceStore,pk=pk)

    if request.method == "POST":
        form = ServiceRename(request.POST, instance=text)
        if form.is_valid():
            doc = form.save(commit=False)
            doc.user = request.user
            doc.save()
            messages.success(request, 'your text have been successfully rename')
            return redirect('service_owner')

    else:
        form = ServiceRename(instance=text)
    return render(request,'rename/service_rename.html',{'form':form})

# NEWS INFORMATION
def news_image_page(request):
    if request.method == 'POST':
        form =NewsImageForm(request.POST,request.FILES)
        if form.is_valid():
            doc = form.save(commit=False)
            doc.user = request.user
            doc.save()
            recent_count = NewsImageStore.objects.filter(show_in_recent=True).count()
            if recent_count > 2:
                oldest_recent = (
                    NewsImageStore.objects.filter(show_in_recent=True)
                    .order_by('upload_at')
                    .first()
                )
                if oldest_recent:
                    oldest_recent.show_in_recent = False
                    oldest_recent.save()
            messages.success(request,'your have successfully uploaded')
            return redirect('mis')
    else:
        form = NewsImageForm()
    return render(request,'upload/news_image.html',{'form':form})

def news_image_delete(request, pk):
    image = get_object_or_404(NewsImageStore,pk=pk,user=request.user)
    image.news_image.delete()
    image.delete()
    return redirect('news_owner')


def news_video_page(request):
    if request.method == 'POST':
        form = NewsVideoForm(request.POST,request.FILES)
        if form.is_valid():
            doc = form.save(commit=False)
            doc.user = request.user
            doc.save()
            recent_count = NewsVideoStore.objects.filter(show_in_recent=True).count()
            if recent_count > 2:
                oldest_recent = (
                    NewsVideoStore.objects.filter(show_in_recent=True)
                    .order_by('upload_at')
                    .first()
                )
                if oldest_recent:
                    oldest_recent.show_in_recent = False
                    oldest_recent.save()
            return redirect('mis')
    else:
        form = NewsVideoForm()
    return render(request,'upload/news_video.html',{'form':form})

def news__video_delete(request, pk):
    video = get_object_or_404(NewsVideoStore,pk=pk, user=request.user)
    video.news_video.delete()
    video.delete()
    return redirect('news_owner')

def news_list(request):
    context = {'news_image':NewsImageStore.objects.all(),
               'news_video':NewsVideoStore.objects.all()
               }
    return render(request,'news.html',{'context':context})

def news_owner(request):
    context = {'news_image':NewsImageStore.objects.filter(user=request.user),
               'news_video':NewsVideoStore.objects.filter(user=request.user)}
    return render(request,'owner/news.html',{'context':context})

def news_image_rename(request,pk):
    image = get_object_or_404(NewsImageStore, pk=pk)

    if request.method == 'POST':
        form = NewsImageRename(request.POST,instance=image)
        if form.is_valid():
            doc = form.save(commit=False)
            doc.user = request.user
            doc.save()
            messages.success(request,'your title have been rename')
            return redirect('news_owner')
    else:
        form = NewsImageRename(instance=image)
    return render(request,'rename/news_image.html',{'form':form})


def news_video_rename(request, pk):
    image = get_object_or_404(NewsVideoStore, pk=pk)

    if request.method == 'POST':
        form = NewsVideoRename(request.POST, instance=image)
        if form.is_valid():
            doc = form.save(commit=False)
            doc.user = request.user
            doc.save()
            messages.success(request, 'your title have been rename')
            return redirect('news_owner')
    else:
        form = NewsImageRename(instance=image)
    return render(request, 'rename/news_video.html', {'form': form})




# slide views information
def slide_page(request):
    if request.method == 'POST':
        form = SlideForm(request.POST,request.FILES)
        if form.is_valid():
            doc = form.save(commit=False)
            doc.user = request.user
            doc.save()
            return redirect('mis')
    else:
        form = SlideForm()
    return render(request,'slide_upload.html',{'form':form})



def slide_delete(request, pk):

    image = get_object_or_404(SlideStore,pk=pk, user=request.user)
    image.slide_image.delete()
    image.delete()
    return redirect('slide_owner')

def slide_list (request):
    image = SlideStore.objects.all()
    news_image = NewsImageStore.objects.all()[:2]
    news_video = NewsVideoStore.objects.all()[:2]
    context = {'image':image,
               'news_image':news_image,
               'news_video':news_video}


    return render(request,'home2.html',{'context':context})




def slide_owner(request):
    files = SlideStore.objects.filter(user=request.user)
    return render(request,'owner/slide.html',{'files':files})



def slide_rename(request, pk):
    image = get_object_or_404(SlideStore, pk=pk)

    if request.method == 'POST':
        form = SlideRename(request.POST, instance=image)
        if form.is_valid():
            doc = form.save(commit=False)
            doc.user = request.user
            doc.save()
            messages.success(request,'your title have been rename successfully')
            return redirect("slide_owner")

    else:
        form = SlideRename(instance=image)
    return render(request,'rename/slide_rename.html',{'form':form})



# TOURIST INFORMATION
def tourist_image_page(request):
    if request.method == 'POST':
        form = TouristImageForm(request.POST,request.FILES)
        if form.is_valid():
            doc = form.save(commit=False)
            doc.user = request.user
            doc.save()
            messages.success(request,'your image have successfully upload')
            return redirect('mis')
    else:
        form = TouristImageForm()
    return render(request,'tour_image.html',{'form':form})

def tourist_video_page(request):
    if request.method == 'POST':
        form = TouristVideoForm(request.POST,request.FILES)
        if form.is_valid():
            doc = form.save(commit=False)
            doc.user = request.user
            doc.save()
            messages.success(request,'your image have successfully upload')
            return redirect('mis')
    else:
        form = TouristVideoForm()
    return render(request,'tour_video.html',{'form':form})

def tourist_image_delete(request,pk):
    file = get_object_or_404(TouristImageStore,pk=pk, user=request.user)
    file.tourist_image.delete()
    file.delete()
    return redirect('tourist_owner')

def tourist_video_delete(request,pk):
    file = get_object_or_404(TouristVideoStore, pk=pk, user=request.user)
    file.tourist_video.delete()
    file.delete()
    return redirect('tourist_owner')


def tourist_list(request):
    context = {'tour_image':TouristImageStore.objects.all(),
               'tour_video':TouristVideoStore.objects.all()}
    return render(request,'tourist.html',{'context':context})

def tourist_owner(request):
    context = {'tour_image':TouristImageStore.objects.filter(user=request.user),
               'tour_video':TouristVideoStore.objects.filter(user=request.user)}
    return render(request,'owner/tourist1.html',{'context':context})


def tourist_image_rename(request, pk):
    image = get_object_or_404(TouristImageStore, pk=pk)

    if request.method == 'POST':
        form = TouristImageRename(request.POST, instance=image)
        if form.is_valid():
            doc = form.save(commit=False)
            doc.user = request.user
            doc.save()
            messages.success(request, 'your title have been rename')
            return redirect('tourist_owner')
    else:
        form = TouristImageRename(instance=image)
    return render(request, 'rename/tour_image.html', {'form': form})

def tourist_video_rename(request, pk):
    image = get_object_or_404(TouristImageStore, pk=pk)

    if request.method == 'POST':
        form = TouristVideoRename(request.POST, instance=image)
        if form.is_valid():
            doc = form.save(commit=False)
            doc.user = request.user
            doc.save()
            messages.success(request, 'your title have been rename')
            return redirect('tourist_owner')
    else:
        form = TouristVideoRename(instance=image)
    return render(request, 'rename/tour_video.html', {'form': form})


# DOCUMENT
def document_page(request):
    if request.method == 'POST':
        form = DocumentForm(request.POST,request.FILES)
        if form.is_valid():
            doc = form.save(commit=False)
            doc.user = request.user
            doc.save()
            messages.success(request,'your image have successfully upload')
            return redirect('mis')
    else:
        form = DocumentForm()
    return render(request,'file_upload.html',{'form':form})

def document_list(request):
    file = DocumentStore.objects.all()
    return render(request,'document.html',{'file':file})

def document_owner(request):
    file = DocumentStore.objects.filter(user=request.user)
    return render(request,'owner/file.html',{'file':file})

def document_delete(request,pk):
    file = get_object_or_404(DocumentStore, pk=pk, user= request.user)
    file.document_upload.delete()
    file.delete()
    return redirect('document_owner')


def document_rename(request, pk):
    image = get_object_or_404(DocumentStore, pk=pk)

    if request.method == 'POST':
        form = DocumentRename(request.POST, instance=image)
        if form.is_valid():
            doc = form.save(commit=False)
            doc.user = request.user
            doc.save()
            messages.success(request, 'your title have been rename')
            return redirect('file_owner')
    else:
        form = DocumentRename(instance=image)
    return render(request, 'rename/document.html', {'form': form})




def view_document(request, pk):
    obj = get_object_or_404(DocumentStore, id=pk)
    file_path = obj.document_upload.path
    ext = os.path.splitext(file_path)[1].lower().replace('.','')

    if ext in ['docx','.docx']:
        content = view_docx(file_path)
    elif ext in ['pptx','.pptx']:
        content =view_pptx(file_path)
    elif ext in ['pdf','.pdf']:
        content = view_pdf(file_path, pk)
    else:
        content ='unsupported file format'
    return render(request, 'open.html', {'content':content})


def view_pdf(file_path,pk ):
    document = get_object_or_404(DocumentStore, pk=pk)
    context = {
        'document': document,
        'file_url': document.document_upload.url,
        'file_path':file_path,
        'file_type':'pdf'
    }
    return context

def view_docx( document):
    doc = DocxDocument(document)

    # Extract paragraphs
    content = []
    for para in doc.paragraphs:
        para_data = {
            'text': para.text,
            'style': para.style.name if para.style else '',
            'alignment': str(para.alignment) if para.alignment else 'LEFT'
        }
        content.append(para_data)

    # Extract tables
    tables_data = []
    for table in doc.tables:
        rows = []
        for row in table.rows:
            cells = [cell.text for cell in row.cells]
            rows.append(cells)
        tables_data.append(rows)

    # Extract images
    images = []
    for rel in doc.part.rels.values():
        if "image" in rel.target_ref:
            try:
                image_data = rel.target_part.blob
                img_base64 = base64.b64encode(image_data).decode('utf-8')
                images.append(f"data:image/png;base64,{img_base64}")
            except Exception:
                pass

    context = {
        'document': document,
        'content': content,
        'tables': tables_data,
        'images': images
    }
    return context

# PPTX viewer
def view_pptx( document):
    prs = Presentation(document)

    slides_data = []
    for slide_num, slide in enumerate(prs.slides):
        slide_data = {'number': slide_num + 1, 'shapes': []}
        for shape in slide.shapes:
            shape_data = {}

            # Text
            if hasattr(shape, "text") and shape.text.strip():
                shape_data['text'] = shape.text.strip()

            # Image
            if hasattr(shape, "image"):
                try:
                    image_bytes = shape.image.blob
                    img_base64 = base64.b64encode(image_bytes).decode('utf-8')
                    shape_data['image'] = f"data:image/png;base64,{img_base64}"
                except Exception:
                    pass

            # Table
            if shape.has_table:
                table_data = []
                for row in shape.table.rows:
                    cells = [cell.text for cell in row.cells]
                    table_data.append(cells)
                shape_data['table'] = table_data

            if shape_data:
                slide_data['shapes'].append(shape_data)

        slides_data.append(slide_data)

    context = {
        'document': document,
        'slides': slides_data
    }
    return context












